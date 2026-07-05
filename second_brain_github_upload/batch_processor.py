import os
import glob
import json
import traceback
from faster_whisper import WhisperModel
import sys

# Forza l'output a usare utf-8 per non far crashare Windows col prompt di Python
try:
    sys.stdout.reconfigure(encoding='utf-8')
except Exception:
    pass

try:
    from pptx import Presentation
except ImportError:
    pass
try:
    import pypdf
except ImportError:
    pass

input_base_dir = r"C:\Users\Schip\.gemini\antigravity\scratch\second_brain_test\input"
output_dir = r"C:\Users\Schip\.gemini\antigravity\scratch\second_brain_test\output"
report_file = r"C:\Users\Schip\.gemini\antigravity\scratch\second_brain_test\processing_report.md"
progress_file = r"C:\Users\Schip\.gemini\antigravity\scratch\second_brain_test\progress.json"

if not os.path.exists(output_dir):
    os.makedirs(output_dir)

def get_unique_filename(filepath, ext):
    parent = os.path.basename(os.path.dirname(filepath))
    basename = os.path.basename(filepath)
    new_name = f"{parent}_{basename}".replace(ext, ".txt")
    # Rigoroso: niente emoji, caratteri strani o nomi esageratamente lunghi
    new_name = "".join(c for c in new_name if c.isalnum() or c in (' ', '.', '_', '-'))
    if len(new_name) > 200:
        new_name = new_name[-200:]
    return os.path.join(output_dir, new_name)

def update_progress(current, total):
    with open(progress_file, "w", encoding="utf-8") as f:
        json.dump({"current": current, "total": total, "percentage": round((current/total)*100, 2)}, f)

def log_report(status, filepath, error_msg=""):
    with open(report_file, "a", encoding="utf-8") as f:
        # Codifichiamo per sicurezza ignorando le emoji così il log non si corrompe
        safe_path = filepath.encode('ascii', 'ignore').decode('ascii')
        if status == "SUCCESS":
            f.write(f"- [x] SUCCESS: {safe_path}\n")
        elif status == "SKIPPED":
            f.write(f"- [/] SKIPPED (Gia' elaborato): {safe_path}\n")
        else:
            safe_err = error_msg.replace('\n', ' ')
            f.write(f"- [ ] ERROR: {safe_path} - Dettaglio: {safe_err}\n")

def process_all():
    print("\n=== AVVIO MEGA-BATCH V2 (Resiliente e Anti-Crash) ===")
    with open(report_file, "w", encoding="utf-8") as f:
        f.write("# Report di Elaborazione Batch (Second Brain)\n\nQuesto file traccia tutti i file processati e gli eventuali errori da recuperare.\n\n")
    
    all_files = []
    
    for root, dirs, files in os.walk(input_base_dir):
        for file in files:
            full_path = os.path.join(root, file)
            ext = full_path.lower()
            if ext.endswith(".mp4") or ext.endswith(".pptx") or ext.endswith(".pdf"):
                all_files.append(full_path)
                
    total_files = len(all_files)
    print(f"Trovati {total_files} file totali da elaborare.")
    
    if total_files == 0:
        return

    # Dividiamo per tipo
    all_mp4 = [f for f in all_files if f.lower().endswith(".mp4")]
    all_pptx = [f for f in all_files if f.lower().endswith(".pptx")]
    all_pdf = [f for f in all_files if f.lower().endswith(".pdf")]
    
    processed_count = 0
    
    # --- PDF ---
    for pdf in all_pdf:
        txt_path = get_unique_filename(pdf, ".pdf")
        if os.path.exists(txt_path):
            log_report("SKIPPED", pdf)
        else:
            try:
                reader = pypdf.PdfReader(pdf)
                text_content = []
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text:
                        text_content.append(f"--- Pagina {i+1} ---\n{text}")
                with open(txt_path, "w", encoding="utf-8") as f:
                    f.write("\n\n".join(text_content))
                log_report("SUCCESS", pdf)
            except Exception as e:
                log_report("ERROR", pdf, str(e))
        processed_count += 1
        update_progress(processed_count, total_files)

    # --- PPTX ---
    for pres in all_pptx:
        txt_path = get_unique_filename(pres, ".pptx")
        if os.path.exists(txt_path):
            log_report("SKIPPED", pres)
        else:
            try:
                prs = Presentation(pres)
                text_content = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    if slide_text:
                        text_content.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
                with open(txt_path, "w", encoding="utf-8") as f:
                    f.write("\n\n".join(text_content))
                log_report("SUCCESS", pres)
            except Exception as e:
                log_report("ERROR", pres, str(e))
        processed_count += 1
        update_progress(processed_count, total_files)

    # --- MP4 ---
    if all_mp4:
        print("Avvio caricamento modello AI Whisper su CPU...")
        model = WhisperModel("base", device="cpu", compute_type="int8")
        for video in all_mp4:
            txt_path = get_unique_filename(video, ".mp4")
            if os.path.exists(txt_path):
                log_report("SKIPPED", video)
            else:
                try:
                    segments, info = model.transcribe(video, beam_size=5, language="it")
                    with open(txt_path, "w", encoding="utf-8") as f:
                        for segment in segments:
                            f.write(f"[{segment.start:.2f}s -> {segment.end:.2f}s] {segment.text}\n")
                    log_report("SUCCESS", video)
                except Exception as e:
                    log_report("ERROR", video, str(e))
            processed_count += 1
            update_progress(processed_count, total_files)

    print("\n=== MEGA-BATCH COMPLETATO CON SUCCESSO! ===")

if __name__ == "__main__":
    process_all()
